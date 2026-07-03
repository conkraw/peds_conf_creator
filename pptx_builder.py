"""PowerPoint export for the Presentation PowerPoint Builder.

The visible deck is intentionally clean: presentation identity appears on the
Title slide only. Speaker-note text is injected into real PowerPoint notes.
"""

from __future__ import annotations

import base64
import html
import io
import math
import re
import zipfile
import posixpath
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from typing import Any, Dict, List, Tuple
from pathlib import Path
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from deck_model import OBJECTIVE_EXAMPLES, slide_output_title, split_nonempty_lines

REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

ET.register_namespace("", REL_NS)
ET.register_namespace("p", P_NS)
ET.register_namespace("r", R_NS)
ET.register_namespace("a", A_NS)

TITLE_BLUE = (31, 78, 121)
LIGHT_BLUE = (234, 242, 250)
PALE_BLUE = (248, 251, 254)
GRAY_TEXT = (80, 80, 80)
BODY_TEXT = (35, 35, 35)
BORDER_BLUE = (185, 205, 225)
WHITE = (255, 255, 255)


# -----------------------------------------------------------------------------
# Basic PowerPoint helpers
# -----------------------------------------------------------------------------


def _rgb(color: Tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def _safe_text(value: Any) -> str:
    return "" if value is None else str(value)


def add_textbox(
    slide,
    text: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int = 22,
    bold: bool = False,
    color: Tuple[int, int, int] = BODY_TEXT,
    align=PP_ALIGN.LEFT,
    fill: Tuple[int, int, int] | None = None,
    margin: float = 0.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _safe_text(text)
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)

    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(fill)
        box.line.color.rgb = _rgb(fill)
    return box


def add_title_bar(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(TITLE_BLUE)
    bar.line.color.rgb = _rgb(TITLE_BLUE)

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.13), Inches(12.3), Inches(0.46))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(23)
    run.font.bold = True
    run.font.color.rgb = _rgb(WHITE)

    if subtitle:
        add_textbox(slide, subtitle, 0.55, 0.78, 12.2, 0.35, 11, False, GRAY_TEXT)


def add_body_lines(slide, lines: List[str], x: float, y: float, w: float, h: float, font_size: int = 22) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.06)

    if not lines:
        lines = ["Add slide content here."]

    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line if line.startswith(("•", "-", "1.", "2.", "3.", "A.", "B.")) else f"• {line}"
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb(BODY_TEXT)
        p.space_after = Pt(7)


def add_section_panel(slide, x: float, y: float, w: float, h: float):
    panel = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(PALE_BLUE)
    panel.line.color.rgb = _rgb(BORDER_BLUE)
    return panel


def estimate_textbox_height(
    text: Any,
    width_inches: float,
    font_size: int,
    min_height: float = 0.28,
    max_height: float = 2.2,
) -> float:
    """Conservative text-height estimate so background boxes grow with text.

    PowerPoint text boxes do not reliably auto-resize when created through
    python-pptx. This estimate prevents the pale-blue title-slide context box
    from being shorter than the Core question / Story arc text.
    """
    clean = _safe_text(text).strip()
    if not clean:
        return 0.0

    chars_per_line = max(18, int(width_inches * (95 / max(font_size, 1))))
    wrapped_lines = 0
    for paragraph in clean.splitlines():
        paragraph = paragraph.strip()
        wrapped_lines += max(1, math.ceil(len(paragraph) / chars_per_line)) if paragraph else 1

    line_height = font_size * 0.020
    estimated = 0.13 + wrapped_lines * line_height
    return max(min_height, min(max_height, estimated))


def add_title_context_panel(
    slide,
    core_question: str,
    story_arc: str,
    x: float,
    y: float,
    w: float,
    max_h: float,
) -> None:
    """Add a Core question / Story arc panel that resizes around its text."""
    core_question = _safe_text(core_question).strip()
    story_arc = _safe_text(story_arc).strip()
    if not core_question and not story_arc:
        return

    content_w = w - 0.60
    core_font = 14
    story_font = 12
    heading_h = 0.24
    gap_after_heading = 0.09
    gap_between_sections = 0.16
    pad_top = 0.18
    pad_bottom = 0.18

    while True:
        core_h = estimate_textbox_height(core_question, content_w, core_font, min_height=0.34, max_height=1.40) if core_question else 0
        story_h = estimate_textbox_height(story_arc, content_w, story_font, min_height=0.34, max_height=1.65) if story_arc else 0
        needed = pad_top + pad_bottom
        if core_question:
            needed += heading_h + gap_after_heading + core_h
        if story_arc:
            needed += (gap_between_sections if core_question else 0) + heading_h + gap_after_heading + story_h
        if needed <= max_h or (core_font <= 11 and story_font <= 9):
            break
        core_font -= 1
        story_font -= 1

    panel_h = min(max_h, max(1.05, needed))
    add_section_panel(slide, x, y, w, panel_h)

    cursor = y + pad_top
    text_x = x + 0.25
    if core_question:
        add_textbox(slide, "Core question", text_x, cursor, content_w, heading_h, 12, True, TITLE_BLUE)
        cursor += heading_h + gap_after_heading
        add_textbox(slide, core_question, text_x, cursor, content_w, core_h, core_font, False, BODY_TEXT, margin=0.04)
        cursor += core_h
    if story_arc:
        if core_question:
            cursor += gap_between_sections
        add_textbox(slide, "Story arc", text_x, cursor, content_w, heading_h, 12, True, TITLE_BLUE)
        cursor += heading_h + gap_after_heading
        add_textbox(slide, story_arc, text_x, cursor, content_w, story_h, story_font, False, BODY_TEXT, margin=0.04)


def get_visual_image(slide_data: Dict[str, Any]) -> Dict[str, str]:
    image = slide_data.get("visual_image", {})
    return image if isinstance(image, dict) else {}


def visual_image_bytes(slide_data: Dict[str, Any]) -> bytes | None:
    encoded = get_visual_image(slide_data).get("data_base64")
    if not encoded:
        return None
    try:
        return base64.b64decode(encoded)
    except Exception:
        return None


def get_uploaded_slide_pptx(slide_data: Dict[str, Any]) -> Dict[str, str]:
    pptx = slide_data.get("uploaded_slide_pptx", {})
    return pptx if isinstance(pptx, dict) else {}


def uploaded_slide_pptx_bytes(slide_data: Dict[str, Any]) -> bytes | None:
    encoded = get_uploaded_slide_pptx(slide_data).get("data_base64")
    if not encoded:
        return None
    try:
        return base64.b64decode(encoded)
    except Exception:
        return None


def render_pptx_first_slide_to_png(pptx_bytes: bytes) -> bytes | None:
    """Render the first slide of an uploaded PPTX to PNG using LibreOffice.

    This avoids directly splicing PowerPoint XML, which can trigger repair or
    read errors in desktop PowerPoint. The rendered slide is then inserted as a
    normal image into the exported deck.
    """
    if not pptx_bytes:
        return None
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            input_path = tmp / "uploaded_slide.pptx"
            input_path.write_bytes(pptx_bytes)
            env = dict(os.environ)
            env["HOME"] = str(tmp)
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "png", "--outdir", str(tmp), str(input_path)],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=90,
                env=env,
            )
            if result.returncode != 0:
                return None
            png_candidates = sorted(tmp.glob("*.png"))
            if not png_candidates:
                return None
            return png_candidates[0].read_bytes()
    except Exception:
        return None


def visual_uses_full_slide(slide_data: Dict[str, Any]) -> bool:
    return bool(slide_data.get("visual_full_slide", False))


def add_image_fit(slide, image_data: bytes, x: float, y: float, w: float, h: float) -> None:
    """Add an uploaded image, preserving aspect ratio inside the given box."""
    try:
        with Image.open(io.BytesIO(image_data)) as img:
            px_w, px_h = img.size
    except Exception:
        add_textbox(slide, "Uploaded visual could not be rendered", x, y, w, 0.5, 12, False, GRAY_TEXT, fill=LIGHT_BLUE)
        return

    if not px_w or not px_h:
        return

    img_ratio = px_w / px_h
    box_ratio = w / h
    if img_ratio >= box_ratio:
        display_w = w
        display_h = w / img_ratio
    else:
        display_h = h
        display_w = h * img_ratio

    left = x + (w - display_w) / 2
    top = y + (h - display_h) / 2
    slide.shapes.add_picture(
        io.BytesIO(image_data),
        Inches(left),
        Inches(top),
        width=Inches(display_w),
        height=Inches(display_h),
    )


def objective_cards_from_slide(slide_data: Dict[str, Any]) -> List[Tuple[str, str]]:
    cards: List[Tuple[str, str]] = []
    defaults = ["Formulate", "Appraise", "Apply"]
    for idx in range(1, 4):
        verb = _safe_text(slide_data.get(f"objective_{idx}_verb", "")).strip()
        text = _safe_text(slide_data.get(f"objective_{idx}_text", "")).strip()
        if verb or text:
            cards.append((verb or defaults[idx - 1], text))
    if cards:
        return cards[:3]

    lines = split_nonempty_lines(slide_data.get("body", "")) or OBJECTIVE_EXAMPLES[:3]
    derived: List[Tuple[str, str]] = []
    fallback_verbs = ["Formulate", "Appraise", "Apply"]
    for idx, line in enumerate(lines[:3], start=1):
        line = _safe_text(line).strip()
        if ":" in line:
            maybe_verb, maybe_text = line.split(":", 1)
            maybe_verb = maybe_verb.strip().title()
            maybe_text = maybe_text.strip()
            if maybe_text:
                derived.append((maybe_verb or fallback_verbs[idx - 1], maybe_text))
                continue
        derived.append((fallback_verbs[idx - 1], line))
    return derived


def estimate_objective_font_size(text: str) -> int:
    words = len(re.findall(r"\b\w+\b", text or ""))
    chars = len((text or "").strip())
    if words > 18 or chars > 110:
        return 17
    if words > 13 or chars > 80:
        return 19
    return 21


def add_objective_card(slide, idx: int, verb: str, text: str, x: float, y: float, w: float, h: float, fill: Tuple[int, int, int]) -> None:
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y - 0.58), Inches(0.60), Inches(0.60))
    circle.fill.solid()
    circle.fill.fore_color.rgb = _rgb(TITLE_BLUE)
    circle.line.color.rgb = _rgb(TITLE_BLUE)
    ctf = circle.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = str(idx)
    cr.font.name = "Aptos"
    cr.font.size = Pt(18)
    cr.font.bold = True
    cr.font.color.rgb = _rgb(WHITE)

    add_textbox(slide, verb or f"Objective {idx}", x + 0.72, y - 0.55, w - 0.2, 0.45, 20, True, TITLE_BLUE)

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + 0.15), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(fill)
    panel.line.color.rgb = _rgb(BORDER_BLUE)
    panel.line.width = Pt(1.0)

    tf = panel.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text or "Add objective text here."
    run.font.name = "Aptos"
    run.font.size = Pt(estimate_objective_font_size(text))
    run.font.color.rgb = _rgb((28, 37, 55))


# -----------------------------------------------------------------------------
# Slide renderers
# -----------------------------------------------------------------------------


def render_title_slide(prs: Presentation, deck: Dict[str, Any], slide_data: Dict[str, Any]) -> None:
    meta = deck.get("metadata", {})
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    title = meta.get("presentation_title") or slide_data.get("title") or "Untitled Presentation"
    subtitle = _safe_text(meta.get("presentation_subtitle", "")).strip()
    presenter = meta.get("presenter") or "Presenter not entered"
    date = meta.get("session_date") or "Date not entered"
    audience = meta.get("audience") or "Audience not entered"
    talk_type = meta.get("presentation_type") or "Presentation type not entered"
    image_data = visual_image_bytes(slide_data)

    if image_data:
        if visual_uses_full_slide(slide_data):
            # Full-title visual mode: use the image as the main slide visual and
            # keep the required title information in a compact readable panel.
            add_image_fit(slide, image_data, 0.35, 0.35, 12.65, 6.80)
            add_section_panel(slide, 0.55, 0.70, 5.85, 2.65)
            title_h = 0.62 if subtitle else 0.95
            add_textbox(slide, title, 0.78, 0.90, 5.35, title_h, 24, True, TITLE_BLUE, fill=PALE_BLUE)
            if subtitle:
                add_textbox(slide, subtitle, 0.82, 1.50, 5.15, 0.32, 13, False, GRAY_TEXT, fill=PALE_BLUE)
            add_textbox(slide, f"{presenter}\n{date}\n{audience}\n{talk_type}", 0.82, 1.92, 5.15, 1.05, 13, False, GRAY_TEXT, fill=PALE_BLUE)
        else:
            # Split title slide: text on the left, optional visual on the right.
            add_textbox(slide, title, 0.70, 0.82, 6.0, 0.92, 30, True, TITLE_BLUE)
            if subtitle:
                add_textbox(slide, subtitle, 0.78, 1.72, 5.7, 0.38, 15, False, GRAY_TEXT)
            add_textbox(slide, f"{presenter}\n{date}\n{audience}\n{talk_type}", 0.78, 2.28, 4.8, 1.35, 17, False, GRAY_TEXT)
            add_section_panel(slide, 7.05, 0.95, 5.55, 5.05)
            add_image_fit(slide, image_data, 7.25, 1.15, 5.15, 4.65)

            core_question = _safe_text(meta.get("core_question", "")).strip()
            story_arc = _safe_text(meta.get("story_arc", "")).strip()
            add_title_context_panel(slide, core_question, story_arc, 0.78, 3.92, 5.95, 2.48)
    else:
        add_textbox(slide, title, 0.75, 0.95, 11.8, 0.90, 34, True, TITLE_BLUE)
        if subtitle:
            add_textbox(slide, subtitle, 0.82, 1.82, 11.2, 0.38, 16, False, GRAY_TEXT)
        add_textbox(slide, f"{presenter}\n{date}\n{audience}\n{talk_type}", 0.80, 2.55, 8.8, 1.15, 18, False, GRAY_TEXT)

        core_question = _safe_text(meta.get("core_question", "")).strip()
        story_arc = _safe_text(meta.get("story_arc", "")).strip()
        add_title_context_panel(slide, core_question, story_arc, 0.80, 4.05, 11.75, 2.25)


def render_objectives_slide(prs: Presentation, deck: Dict[str, Any], slide_data: Dict[str, Any], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    image_data = visual_image_bytes(slide_data)
    if image_data and visual_uses_full_slide(slide_data):
        add_textbox(slide, slide_output_title(deck, slide_data, index), 0.30, 0.18, 4.2, 0.48, 24, True, TITLE_BLUE)
        add_image_fit(slide, image_data, 0.35, 0.80, 12.60, 6.30)
        return

    title = slide_output_title(deck, slide_data, index)
    intro = _safe_text(slide_data.get("objectives_intro", "")).strip() or "By the end of this session, residents should be able to:"
    takeaway = _safe_text(slide_data.get("objectives_takeaway", "")).strip()

    # More visual custom layout to match the reference style.
    add_textbox(slide, title, 0.30, 0.16, 5.8, 0.50, 24, True, TITLE_BLUE)
    add_textbox(slide, intro, 0.30, 0.70, 12.3, 0.24, 12, False, GRAY_TEXT)

    divider = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.30), Inches(1.03), Inches(12.45), Inches(0.015))
    divider.fill.solid()
    divider.fill.fore_color.rgb = _rgb(BORDER_BLUE)
    divider.line.color.rgb = _rgb(BORDER_BLUE)

    cards = objective_cards_from_slide(slide_data)
    while len(cards) < 3:
        cards.append((f"Objective {len(cards)+1}", "Add objective text here."))

    fills = [(235, 241, 247), (228, 238, 223), (243, 232, 193)]
    x_positions = [0.54, 4.95, 9.36]
    y = 2.05
    w = 3.80
    h = 2.60
    for idx_card, (x, card) in enumerate(zip(x_positions, cards[:3]), start=1):
        verb, sentence = card
        add_objective_card(slide, idx_card, verb, sentence, x, y, w, h, fills[(idx_card - 1) % len(fills)])

    if takeaway:
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(5.77), Inches(11.70), Inches(0.78))
        panel.fill.solid()
        panel.fill.fore_color.rgb = _rgb(WHITE)
        panel.line.color.rgb = _rgb(BORDER_BLUE)
        panel.line.width = Pt(1.0)
        tf = panel.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = takeaway
        run.font.name = "Aptos"
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = _rgb((25, 35, 52))


def render_disclosures_slide(prs: Presentation, deck: Dict[str, Any], slide_data: Dict[str, Any], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, slide_output_title(deck, slide_data, index))
    disclosures = split_nonempty_lines(slide_data.get("body", "")) or ["I have no relevant financial or non-financial disclosures."]
    add_body_lines(slide, disclosures, 0.95, 1.55, 11.0, 4.2, 23)


def render_standard_slide(prs: Presentation, deck: Dict[str, Any], slide_data: Dict[str, Any], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide_output_title(deck, slide_data, index)
    subtitle = _safe_text(slide_data.get("subtitle", "")).strip()
    add_title_bar(slide, title, subtitle)

    body_lines = split_nonempty_lines(slide_data.get("body", ""))
    discussion = _safe_text(slide_data.get("discussion_prompt", "")).strip()
    image_data = visual_image_bytes(slide_data)

    if image_data:
        if visual_uses_full_slide(slide_data):
            # Full visual mode keeps the standard title bar but gives the image
            # essentially the entire slide body. Text and discussion remain in
            # the mentor document and speaker notes, not on this slide.
            add_image_fit(slide, image_data, 0.55, 0.95, 12.25, 6.10)
        else:
            # Default visual mode: image takes about half of the slide.
            text_x = 0.75
            text_w = 5.55
            image_x = 6.65
            image_w = 5.95
            top_y = 1.22
            content_h = 5.45

            if discussion:
                add_body_lines(slide, body_lines, text_x, top_y, text_w, 3.85, 20)
                add_textbox(slide, "Discussion prompt", text_x, 5.18, text_w, 0.28, 13, True, TITLE_BLUE)
                add_textbox(slide, discussion, text_x, 5.50, text_w, 0.92, 12, False, BODY_TEXT, fill=PALE_BLUE)
            else:
                add_body_lines(slide, body_lines, text_x, top_y, text_w, content_h, 21)

            add_image_fit(slide, image_data, image_x, top_y, image_w, content_h)
    else:
        if discussion:
            add_body_lines(slide, body_lines, 0.85, 1.35, 11.5, 4.45, 22)
            add_textbox(slide, "Discussion prompt", 0.95, 5.95, 2.3, 0.28, 13, True, TITLE_BLUE)
            add_textbox(slide, discussion, 0.95, 6.24, 11.2, 0.50, 12, False, BODY_TEXT, fill=PALE_BLUE)
        else:
            add_body_lines(slide, body_lines, 0.85, 1.35, 11.5, 5.25, 22)


def _rels_path_for(part_path: str) -> str:
    folder = posixpath.dirname(part_path)
    name = posixpath.basename(part_path)
    return posixpath.join(folder, "_rels", f"{name}.rels")


def _resolve_rel_target(part_path: str, target: str) -> str:
    return posixpath.normpath(posixpath.join(posixpath.dirname(part_path), target))


def _relative_target(from_part: str, to_part: str) -> str:
    return posixpath.relpath(to_part, posixpath.dirname(from_part))


def _content_type_maps(files: Dict[str, bytes]) -> Tuple[Dict[str, str], Dict[str, str]]:
    root = ET.fromstring(files["[Content_Types].xml"])
    overrides = {}
    defaults = {}
    for node in root.findall(f"{{{CT_NS}}}Override"):
        overrides[node.get("PartName", "").lstrip("/")] = node.get("ContentType", "")
    for node in root.findall(f"{{{CT_NS}}}Default"):
        defaults[node.get("Extension", "").lower()] = node.get("ContentType", "")
    return overrides, defaults


def _next_unique_part_path(dest_files: Dict[str, bytes], desired_path: str) -> str:
    if desired_path not in dest_files:
        return desired_path
    base_dir = posixpath.dirname(desired_path)
    stem, ext = posixpath.splitext(posixpath.basename(desired_path))
    counter = 1
    while True:
        candidate = posixpath.join(base_dir, f"{stem}_import{counter}{ext}")
        if candidate not in dest_files:
            return candidate
        counter += 1


def _copy_part_recursive(
    src_files: Dict[str, bytes],
    dest_files: Dict[str, bytes],
    src_overrides: Dict[str, str],
    src_defaults: Dict[str, str],
    src_part_path: str,
    memo: Dict[str, str],
    forced_dest_part: str | None = None,
) -> str:
    if forced_dest_part is None and src_part_path in memo:
        return memo[src_part_path]

    desired_dest = forced_dest_part or src_part_path
    src_bytes = src_files[src_part_path]
    if desired_dest in dest_files and dest_files[desired_dest] != src_bytes and forced_dest_part is None:
        desired_dest = _next_unique_part_path(dest_files, desired_dest)

    memo[src_part_path] = desired_dest
    dest_files[desired_dest] = src_bytes

    rels_src_path = _rels_path_for(src_part_path)
    rels_dest_path = _rels_path_for(desired_dest)
    if rels_src_path in src_files:
        rel_root = ET.fromstring(src_files[rels_src_path])
        for rel in list(rel_root.findall(f"{{{REL_NS}}}Relationship")):
            rel_type = rel.get("Type", "")
            if rel_type.endswith(("/notesSlide", "/comments", "/commentAuthors")):
                rel_root.remove(rel)
                continue
            if rel.get("TargetMode") == "External":
                continue
            target = rel.get("Target", "")
            if not target:
                continue
            src_target_part = _resolve_rel_target(src_part_path, target)
            if src_target_part not in src_files:
                continue
            dest_target_part = _copy_part_recursive(src_files, dest_files, src_overrides, src_defaults, src_target_part, memo)
            rel.set("Target", _relative_target(desired_dest, dest_target_part))
        dest_files[rels_dest_path] = ET.tostring(rel_root, encoding="utf-8", xml_declaration=True)
    else:
        dest_files.pop(rels_dest_path, None)

    part_name = f"/{desired_dest}"
    if part_name in src_overrides:
        dest_files["[Content_Types].xml"] = _add_content_type_override(dest_files["[Content_Types].xml"], part_name, src_overrides[part_name])
    else:
        ext = posixpath.splitext(desired_dest)[1].lstrip(".").lower()
        if ext and ext in src_defaults:
            dest_files["[Content_Types].xml"] = _add_content_type_default(dest_files["[Content_Types].xml"], ext, src_defaults[ext])
    return desired_dest



def _ensure_slide_master_registered(dest_files: Dict[str, bytes], slide_number: int) -> None:
    """Register the imported slide's master in presentation.xml if needed.

    PowerPoint can repair imported slides when a slide points to an imported
    slideLayout whose slideMaster is present in the package but not listed in
    ppt/presentation.xml. This makes the master relationship explicit.
    """
    slide_part = f"ppt/slides/slide{slide_number}.xml"
    slide_rels_path = _rels_path_for(slide_part)
    if slide_rels_path not in dest_files:
        return

    try:
        slide_rels = ET.fromstring(dest_files[slide_rels_path])
        layout_part = None
        for rel in slide_rels.findall(f"{{{REL_NS}}}Relationship"):
            if rel.get("Type", "").endswith("/slideLayout"):
                layout_part = _resolve_rel_target(slide_part, rel.get("Target", ""))
                break
        if not layout_part:
            return

        layout_rels_path = _rels_path_for(layout_part)
        if layout_rels_path not in dest_files:
            return
        layout_rels = ET.fromstring(dest_files[layout_rels_path])
        master_part = None
        for rel in layout_rels.findall(f"{{{REL_NS}}}Relationship"):
            if rel.get("Type", "").endswith("/slideMaster"):
                master_part = _resolve_rel_target(layout_part, rel.get("Target", ""))
                break
        if not master_part or master_part not in dest_files:
            return

        pres_rels_path = "ppt/_rels/presentation.xml.rels"
        pres_path = "ppt/presentation.xml"
        pres_rels = ET.fromstring(dest_files[pres_rels_path])
        existing_rid = None
        for rel in pres_rels.findall(f"{{{REL_NS}}}Relationship"):
            if rel.get("Type", "").endswith("/slideMaster"):
                rel_target = _resolve_rel_target("ppt/presentation.xml", rel.get("Target", ""))
                if rel_target == master_part:
                    existing_rid = rel.get("Id")
                    break

        if existing_rid is None:
            existing_rid = _next_rid(pres_rels)
            ET.SubElement(
                pres_rels,
                f"{{{REL_NS}}}Relationship",
                Id=existing_rid,
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
                Target=_relative_target("ppt/presentation.xml", master_part),
            )
            dest_files[pres_rels_path] = ET.tostring(pres_rels, encoding="utf-8", xml_declaration=True)

        pres = ET.fromstring(dest_files[pres_path])
        ns = {"p": P_NS, "r": R_NS}
        master_list = pres.find("p:sldMasterIdLst", ns)
        if master_list is None:
            master_list = ET.Element(f"{{{P_NS}}}sldMasterIdLst")
            pres.insert(0, master_list)

        for master_id in master_list.findall(f"{{{P_NS}}}sldMasterId"):
            if master_id.get(f"{{{R_NS}}}id") == existing_rid:
                return

        used_ids = []
        for master_id in master_list.findall(f"{{{P_NS}}}sldMasterId"):
            try:
                used_ids.append(int(master_id.get("id", "0")))
            except Exception:
                pass
        new_id = max(used_ids or [2147483647]) + 1
        master_id = ET.SubElement(master_list, f"{{{P_NS}}}sldMasterId")
        master_id.set("id", str(new_id))
        master_id.set(f"{{{R_NS}}}id", existing_rid)
        dest_files[pres_path] = ET.tostring(pres, encoding="utf-8", xml_declaration=True)
    except Exception:
        return

def _source_first_slide_part(src_files: Dict[str, bytes]) -> str:
    try:
        pres = ET.fromstring(src_files["ppt/presentation.xml"])
        ns = {"p": P_NS, "r": R_NS}
        slide_id = pres.find("p:sldIdLst/p:sldId", ns)
        if slide_id is not None:
            rel_id = slide_id.get(f"{{{R_NS}}}id")
            rels = ET.fromstring(src_files["ppt/_rels/presentation.xml.rels"])
            for rel in rels.findall(f"{{{REL_NS}}}Relationship"):
                if rel.get("Id") == rel_id:
                    return _resolve_rel_target("ppt/presentation.xml", rel.get("Target", "slides/slide1.xml"))
    except Exception:
        pass
    return "ppt/slides/slide1.xml"


def render_uploaded_pptx_replacement_slide(prs: Presentation, deck: Dict[str, Any], slide_data: Dict[str, Any], index: int) -> None:
    """Render an uploaded PPTX's first slide as a full-slide image.

    The key design choice is intentional: imported PPTX slides are rasterized to
    a PNG and inserted as a normal image. This is far more reliable than XML
    splicing and prevents desktop PowerPoint repair/read errors.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pptx_bytes = uploaded_slide_pptx_bytes(slide_data)
    image_bytes = render_pptx_first_slide_to_png(pptx_bytes or b"")
    if image_bytes:
        add_image_fit(slide, image_bytes, 0.0, 0.0, 13.333, 7.50)
        return

    add_title_bar(slide, slide_output_title(deck, slide_data, index))
    info = get_uploaded_slide_pptx(slide_data)
    filename = info.get("filename", "uploaded slide.pptx")
    add_section_panel(slide, 0.85, 1.50, 11.65, 3.75)
    add_textbox(slide, "PPTX slide could not be rendered", 1.15, 1.85, 11.0, 0.45, 24, True, TITLE_BLUE)
    add_textbox(
        slide,
        f"The uploaded PPTX was stored, but the app could not render its first slide as an image.\n\nFile: {filename}\n\nOn Streamlit Cloud, make sure packages.txt includes libreoffice. You can also export the slide as PNG and upload that image.",
        1.15,
        2.45,
        10.95,
        1.95,
        17,
        False,
        BODY_TEXT,
    )

def replace_uploaded_pptx_slides(pptx_bytes: bytes, deck: Dict[str, Any]) -> bytes:
    with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zin:
        dest_files = {name: zin.read(name) for name in zin.namelist()}

    for slide_number, slide_data in enumerate(deck.get("slides", []), start=1):
        replacement_bytes = uploaded_slide_pptx_bytes(slide_data)
        if not replacement_bytes:
            continue
        try:
            with zipfile.ZipFile(io.BytesIO(replacement_bytes), "r") as src_zip:
                src_files = {name: src_zip.read(name) for name in src_zip.namelist()}
            src_overrides, src_defaults = _content_type_maps(src_files)
            src_slide_part = _source_first_slide_part(src_files)
            if src_slide_part not in src_files:
                continue
            memo: Dict[str, str] = {}
            _copy_part_recursive(
                src_files,
                dest_files,
                src_overrides,
                src_defaults,
                src_slide_part,
                memo,
                forced_dest_part=f"ppt/slides/slide{slide_number}.xml",
            )
            _ensure_slide_master_registered(dest_files, slide_number)
        except Exception:
            continue

    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in dest_files.items():
            zout.writestr(name, data)
    return output.getvalue()


def build_pptx(deck: Dict[str, Any]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    speaker_notes: List[str] = []
    output_index = 0
    for slide_data in deck.get("slides", []):
        output_index += 1
        kind = slide_data.get("slide_kind")
        role = slide_data.get("role")
        if uploaded_slide_pptx_bytes(slide_data):
            render_uploaded_pptx_replacement_slide(prs, deck, slide_data, output_index)
        elif kind == "title" or role == "Title":
            render_title_slide(prs, deck, slide_data)
        elif kind == "objectives" or role == "Objectives":
            render_objectives_slide(prs, deck, slide_data, output_index)
        elif kind == "disclosures" or role == "Disclosures":
            render_disclosures_slide(prs, deck, slide_data, output_index)
        else:
            render_standard_slide(prs, deck, slide_data, output_index)
        speaker_notes.append(_safe_text(slide_data.get("speaker_notes", "")))

    buffer = io.BytesIO()
    prs.save(buffer)
    pptx_bytes = buffer.getvalue()
    return add_speaker_notes_to_pptx(pptx_bytes, speaker_notes)


# -----------------------------------------------------------------------------
# Speaker-notes Office XML injection
# -----------------------------------------------------------------------------


def _next_rid(root: ET.Element) -> str:
    nums: List[int] = []
    for rel in root.findall(f"{{{REL_NS}}}Relationship"):
        rid = rel.get("Id", "")
        match = re.match(r"rId(\d+)$", rid)
        if match:
            nums.append(int(match.group(1)))
    return f"rId{max(nums or [0]) + 1}"


def _add_content_type_override(content_types_xml: bytes, part_name: str, content_type: str) -> bytes:
    ET.register_namespace("", CT_NS)
    root = ET.fromstring(content_types_xml)
    for override in root.findall(f"{{{CT_NS}}}Override"):
        if override.get("PartName") == part_name:
            override.set("ContentType", content_type)
            return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    ET.SubElement(root, f"{{{CT_NS}}}Override", PartName=part_name, ContentType=content_type)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _add_content_type_default(content_types_xml: bytes, extension: str, content_type: str) -> bytes:
    ET.register_namespace("", CT_NS)
    root = ET.fromstring(content_types_xml)
    for default in root.findall(f"{{{CT_NS}}}Default"):
        if default.get("Extension", "").lower() == extension.lower().lstrip("."):
            return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    ET.SubElement(root, f"{{{CT_NS}}}Default", Extension=extension.lstrip("."), ContentType=content_type)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _notes_master_xml() -> bytes:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notesMaster xmlns:a="{A_NS}" xmlns:r="{R_NS}" xmlns:p="{P_NS}">
  <p:cSld>
    <p:bg><p:bgRef idx="1001"><a:schemeClr val="bg1"/></p:bgRef></p:bg>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg" idx="2"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="914400"/><a:ext cx="5486400" cy="3086100"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln w="12700"><a:solidFill><a:prstClr val="black"/></a:solidFill></a:ln></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
      </p:sp>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="4400550"/><a:ext cx="5486400" cy="3600450"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="en-US"/><a:t>Speaker notes</a:t></a:r></a:p></p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHLink"/>
</p:notesMaster>'''.encode("utf-8")


def _notes_slide_xml(notes: str, slide_number: int) -> bytes:
    lines = notes.splitlines() or [""]
    paragraphs = []
    for line in lines:
        escaped = html.escape(line, quote=False)
        paragraphs.append(
            f'<a:p><a:r><a:rPr lang="en-US" dirty="0"/><a:t>{escaped}</a:t></a:r><a:endParaRPr lang="en-US" dirty="0"/></a:p>'
        )
    paragraphs_xml = "".join(paragraphs)
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="{A_NS}" xmlns:r="{R_NS}" xmlns:p="{P_NS}">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr><p:spPr/></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/>{paragraphs_xml}</p:txBody></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="4" name="Slide Number Placeholder 3"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldNum" sz="quarter" idx="10"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:fld id="{{F7021451-1387-4CA6-816F-3879F97B5CBC}}" type="slidenum"><a:rPr lang="en-US"/><a:t>{slide_number}</a:t></a:fld><a:endParaRPr lang="en-US"/></a:p></p:txBody></p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>'''.encode("utf-8")


def _notes_slide_rels_xml(slide_number: int) -> bytes:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL_NS}">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{slide_number}.xml"/>
</Relationships>'''.encode("utf-8")


def add_speaker_notes_to_pptx(pptx_bytes: bytes, notes_by_slide: List[str]) -> bytes:
    """Add real PowerPoint speaker notes to a python-pptx-generated deck."""
    if not any((note or "").strip() for note in notes_by_slide):
        return pptx_bytes

    with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    files["[Content_Types].xml"] = _add_content_type_override(
        files["[Content_Types].xml"],
        "/ppt/notesMasters/notesMaster1.xml",
        "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml",
    )

    presentation_rels_path = "ppt/_rels/presentation.xml.rels"
    pres_rels = ET.fromstring(files[presentation_rels_path])
    notes_master_rid = None
    for rel in pres_rels.findall(f"{{{REL_NS}}}Relationship"):
        if rel.get("Type", "").endswith("/notesMaster"):
            notes_master_rid = rel.get("Id")
            break
    if not notes_master_rid:
        notes_master_rid = _next_rid(pres_rels)
        ET.SubElement(
            pres_rels,
            f"{{{REL_NS}}}Relationship",
            Id=notes_master_rid,
            Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster",
            Target="notesMasters/notesMaster1.xml",
        )
        files[presentation_rels_path] = ET.tostring(pres_rels, encoding="utf-8", xml_declaration=True)

    presentation_path = "ppt/presentation.xml"
    pres = ET.fromstring(files[presentation_path])
    ns = {"p": P_NS, "r": R_NS}
    if pres.find("p:notesMasterIdLst", ns) is None:
        notes_master_id_list = ET.Element(f"{{{P_NS}}}notesMasterIdLst")
        notes_master_id = ET.SubElement(notes_master_id_list, f"{{{P_NS}}}notesMasterId")
        notes_master_id.set(f"{{{R_NS}}}id", notes_master_rid)
        insert_at = len(list(pres))
        for i, child in enumerate(list(pres)):
            if child.tag == f"{{{P_NS}}}sldSz":
                insert_at = i
                break
        pres.insert(insert_at, notes_master_id_list)
        files[presentation_path] = ET.tostring(pres, encoding="utf-8", xml_declaration=True)

    files["ppt/notesMasters/notesMaster1.xml"] = _notes_master_xml()
    files["ppt/notesMasters/_rels/notesMaster1.xml.rels"] = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL_NS}"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/></Relationships>'''.encode("utf-8")

    for slide_number, note in enumerate(notes_by_slide, start=1):
        if not (note or "").strip():
            continue
        slide_rels_path = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
        if slide_rels_path not in files:
            continue
        slide_rels = ET.fromstring(files[slide_rels_path])
        already_linked = any(rel.get("Type", "").endswith("/notesSlide") for rel in slide_rels.findall(f"{{{REL_NS}}}Relationship"))
        if not already_linked:
            rid = _next_rid(slide_rels)
            ET.SubElement(
                slide_rels,
                f"{{{REL_NS}}}Relationship",
                Id=rid,
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
                Target=f"../notesSlides/notesSlide{slide_number}.xml",
            )
            files[slide_rels_path] = ET.tostring(slide_rels, encoding="utf-8", xml_declaration=True)

        files[f"ppt/notesSlides/notesSlide{slide_number}.xml"] = _notes_slide_xml(note, slide_number)
        files[f"ppt/notesSlides/_rels/notesSlide{slide_number}.xml.rels"] = _notes_slide_rels_xml(slide_number)
        files["[Content_Types].xml"] = _add_content_type_override(
            files["[Content_Types].xml"],
            f"/ppt/notesSlides/notesSlide{slide_number}.xml",
            "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
        )

    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)
    return output.getvalue()

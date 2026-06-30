"""
Pediatric Residency Presentation Builder
---------------------------------------
A Streamlit app that standardizes educational/study-review PowerPoint creation,
keeps the talk narrative-driven, supports mentor review, writes real PowerPoint
speaker notes, and can save drafts + final PPTX files to GitHub.

Expected Streamlit secrets for GitHub saving:
GITHUB_TOKEN = "ghp_..."
GITHUB_REPO = "your_username/your_repo"
GITHUB_BRANCH = "main"
GITHUB_FOLDER = "presentation_archive"
"""

from __future__ import annotations

import base64
import copy
import datetime as dt
import html
import io
import json
import re
import uuid
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple

import requests
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


APP_TITLE = "Pediatric Residency Presentation Builder"
APP_VERSION = "2026.06.30"
ARCHIVE_JSON_NAME = "draft.json"
ARCHIVE_PPTX_NAME = "presentation.pptx"

# PowerPoint XML namespaces for real speaker-note injection.
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
ET.register_namespace("", REL_NS)
ET.register_namespace("p", P_NS)
ET.register_namespace("r", R_NS)
ET.register_namespace("a", A_NS)


# -----------------------------------------------------------------------------
# Data model
# -----------------------------------------------------------------------------


def new_slide(
    role: str,
    title: str,
    prompt: str = "",
    slide_type: str = "Standard",
    required: bool = False,
) -> Dict[str, Any]:
    return {
        "id": str(uuid.uuid4()),
        "role": role,
        "title": title,
        "subtitle": "",
        "prompt": prompt,
        "body": "",
        "visual_plan": "",
        "discussion_prompt": "",
        "speaker_notes": "",
        "slide_type": slide_type,
        "include": True,
        "required": required,
    }


def default_deck() -> Dict[str, Any]:
    today = dt.date.today().isoformat()
    return {
        "app_version": APP_VERSION,
        "metadata": {
            "presentation_title": "",
            "presenter": "",
            "session_date": today,
            "audience": "Pediatric residents",
            "presentation_type": "Educational Topic",
            "core_question": "",
            "story_arc": "",
        },
        "mentor_review": {
            "mentor_name": "",
            "mentor_email": "",
            "review_status": "Not sent",
            "requested_review_date": "",
            "review_completed_date": "",
            "mentor_feedback": "",
            "mentor_approval_statement": "",
            "include_mentor_review_slide": False,
            "checklist": {
                "objectives_use_bloom": False,
                "slide_titles_tell_story": False,
                "no_data_dump": False,
                "speaker_notes_complete": False,
                "take_home_points_clear": False,
            },
        },
        "slides": [
            new_slide("Title", "", "Name the talk clearly and state the clinical/research problem.", "Title", True),
            new_slide("Objectives", "Objectives", "Use measurable verbs from Bloom's taxonomy.", "Objectives", True),
            new_slide("Disclosures", "Disclosures", "State relevant financial/non-financial disclosures or no relevant disclosures.", "Standard", True),
            new_slide("Introduction", "Why this matters", "Orient the audience: why should residents care in the first 60 seconds?", "Standard", True),
            new_slide("Story setup", "The question we need to answer", "Create tension: clinical dilemma, knowledge gap, or study question.", "Story", True),
            new_slide("Core content", "Key idea 1", "Teach the first major concept, method, or result.", "Story", False),
            new_slide("Core content", "Key idea 2", "Teach the second major concept, method, or result.", "Story", False),
            new_slide("Application", "How this changes our thinking", "Return to patient care, interpretation, or practical decision-making.", "Story", False),
            new_slide("Take-home", "Take-home points", "End with 2–3 actionable points.", "Take-home", True),
        ],
    }


BLOOM_HELPER = {
    "Remember": "define, list, identify, name, recall",
    "Understand": "describe, summarize, explain, classify, compare",
    "Apply": "use, demonstrate, calculate, choose, implement",
    "Analyze": "differentiate, organize, interpret, examine, contrast",
    "Evaluate": "appraise, justify, critique, prioritize, defend",
    "Create": "design, formulate, develop, construct, propose",
}

OBJECTIVE_EXAMPLES = [
    "Describe the clinical problem and why it matters for pediatric practice.",
    "Differentiate the key diagnostic or management options using case data.",
    "Appraise the strength and limitations of the evidence being presented.",
    "Apply the take-home points to a realistic patient-care decision.",
]

CHECKLIST_LABELS = {
    "objectives_use_bloom": "Objectives use Bloom-style measurable verbs",
    "slide_titles_tell_story": "Slide titles tell the story rather than label topics only",
    "no_data_dump": "Slides avoid data dumping and each slide has one job",
    "speaker_notes_complete": "Speaker notes are complete enough for rehearsal/review",
    "take_home_points_clear": "Final take-home points are clear and actionable",
}


# -----------------------------------------------------------------------------
# General helpers
# -----------------------------------------------------------------------------


def initialize_state() -> None:
    if "deck" not in st.session_state:
        st.session_state.deck = default_deck()
    if "selected_slide_id" not in st.session_state:
        st.session_state.selected_slide_id = st.session_state.deck["slides"][0]["id"]


def sanitize_filename(value: str, default: str = "presentation") -> str:
    value = (value or "").strip() or default
    value = re.sub(r"[^A-Za-z0-9._ -]+", "", value)
    value = re.sub(r"\s+", "_", value)
    return value[:90].strip("._-") or default


def split_nonempty_lines(text: str) -> List[str]:
    return [line.strip() for line in (text or "").splitlines() if line.strip()]


def short_label(text: str, max_len: int = 34) -> str:
    text = (text or "").strip()
    return text if len(text) <= max_len else text[: max_len - 1].rstrip() + "…"


def get_selected_slide(deck: Dict[str, Any]) -> Dict[str, Any]:
    slide_id = st.session_state.selected_slide_id
    for slide in deck["slides"]:
        if slide["id"] == slide_id:
            return slide
    st.session_state.selected_slide_id = deck["slides"][0]["id"]
    return deck["slides"][0]


def slide_output_title(deck: Dict[str, Any], slide: Dict[str, Any], index: int) -> str:
    title = (slide.get("title") or "").strip()
    if title:
        return title
    if slide.get("role") == "Title":
        return deck["metadata"].get("presentation_title") or "Untitled Presentation"
    return f"Story slide {index}"


def make_archive_slug(deck: Dict[str, Any]) -> str:
    meta = deck["metadata"]
    date = sanitize_filename(meta.get("session_date") or dt.date.today().isoformat(), "date")
    title = sanitize_filename(meta.get("presentation_title") or "untitled_presentation")
    presenter = sanitize_filename(meta.get("presenter") or "presenter")
    return f"{date}_{presenter}_{title}"


def to_json_bytes(deck: Dict[str, Any]) -> bytes:
    payload = copy.deepcopy(deck)
    payload["saved_at"] = dt.datetime.now().isoformat(timespec="seconds")
    return json.dumps(payload, indent=2, ensure_ascii=False).encode("utf-8")


def load_deck_from_json(uploaded_bytes: bytes) -> Dict[str, Any]:
    loaded = json.loads(uploaded_bytes.decode("utf-8"))
    # Keep backward compatibility if future fields were missing.
    base = default_deck()
    base.update(loaded)
    base.setdefault("metadata", default_deck()["metadata"])
    base.setdefault("mentor_review", default_deck()["mentor_review"])
    base.setdefault("slides", default_deck()["slides"])
    return base


# -----------------------------------------------------------------------------
# GitHub persistence
# -----------------------------------------------------------------------------


@dataclass
class GitHubConfig:
    token: str
    repo: str
    branch: str
    folder: str


def github_config_from_secrets() -> Optional[GitHubConfig]:
    token = st.secrets.get("GITHUB_TOKEN", "") if hasattr(st, "secrets") else ""
    repo = st.secrets.get("GITHUB_REPO", "") if hasattr(st, "secrets") else ""
    branch = st.secrets.get("GITHUB_BRANCH", "main") if hasattr(st, "secrets") else "main"
    folder = st.secrets.get("GITHUB_FOLDER", "presentation_archive") if hasattr(st, "secrets") else "presentation_archive"
    if not token or not repo:
        return None
    return GitHubConfig(token=token, repo=repo, branch=branch, folder=folder.strip("/"))


def github_headers(config: GitHubConfig) -> Dict[str, str]:
    return {
        "Authorization": f"Bearer {config.token}",
        "Accept": "application/vnd.github+json",
        "X-GitHub-Api-Version": "2022-11-28",
    }


def github_api_url(config: GitHubConfig, path: str) -> str:
    safe_path = "/".join(part.strip("/") for part in path.split("/") if part.strip("/"))
    return f"https://api.github.com/repos/{config.repo}/contents/{safe_path}"


def github_get_file_sha(config: GitHubConfig, path: str) -> Optional[str]:
    response = requests.get(
        github_api_url(config, path),
        headers=github_headers(config),
        params={"ref": config.branch},
        timeout=30,
    )
    if response.status_code == 404:
        return None
    response.raise_for_status()
    data = response.json()
    return data.get("sha")


def github_upsert_bytes(config: GitHubConfig, path: str, content: bytes, message: str) -> None:
    sha = github_get_file_sha(config, path)
    payload: Dict[str, Any] = {
        "message": message,
        "content": base64.b64encode(content).decode("ascii"),
        "branch": config.branch,
    }
    if sha:
        payload["sha"] = sha
    response = requests.put(
        github_api_url(config, path),
        headers=github_headers(config),
        json=payload,
        timeout=60,
    )
    response.raise_for_status()


def github_list_json_drafts(config: GitHubConfig) -> List[Dict[str, Any]]:
    response = requests.get(
        github_api_url(config, config.folder),
        headers=github_headers(config),
        params={"ref": config.branch},
        timeout=30,
    )
    if response.status_code == 404:
        return []
    response.raise_for_status()
    items = response.json()
    if not isinstance(items, list):
        return []

    drafts: List[Dict[str, Any]] = []
    for item in items:
        if item.get("type") != "dir":
            continue
        folder_path = item.get("path", "")
        draft_path = f"{folder_path}/{ARCHIVE_JSON_NAME}"
        draft_resp = requests.get(
            github_api_url(config, draft_path),
            headers=github_headers(config),
            params={"ref": config.branch},
            timeout=30,
        )
        if draft_resp.status_code == 200:
            drafts.append({"name": item.get("name"), "path": draft_path})
    return sorted(drafts, key=lambda item: item["name"], reverse=True)


def github_load_json(config: GitHubConfig, path: str) -> Dict[str, Any]:
    response = requests.get(
        github_api_url(config, path),
        headers=github_headers(config),
        params={"ref": config.branch},
        timeout=30,
    )
    response.raise_for_status()
    data = response.json()
    raw = base64.b64decode(data["content"])
    return load_deck_from_json(raw)


# -----------------------------------------------------------------------------
# PowerPoint generation
# -----------------------------------------------------------------------------


def add_textbox(slide, text: str, x, y, w, h, font_size=24, bold=False, color=(30, 30, 30)):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_title_bar(slide, title: str, subtitle: str = "") -> None:
    # Deep blue title band with simple academic feel.
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(31, 78, 121)
    bar.line.color.rgb = RGBColor(31, 78, 121)
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.13), Inches(12.3), Inches(0.44))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(23)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.55), Inches(0.78), Inches(12.2), Inches(0.32), 11, False, (85, 85, 85))


def add_body_lines(slide, lines: List[str], x, y, w, h, font_size=22) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    if not lines:
        lines = ["Add slide content here."]
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)
        p.space_after = Pt(7)


def render_standard_slide(prs: Presentation, deck: Dict[str, Any], slide_data: Dict[str, Any], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide_output_title(deck, slide_data, index)
    subtitle = slide_data.get("subtitle", "")
    add_title_bar(slide, title, subtitle)

    body_lines = split_nonempty_lines(slide_data.get("body", ""))
    visual_plan = (slide_data.get("visual_plan") or "").strip()
    discussion = (slide_data.get("discussion_prompt") or "").strip()

    # Main content box.
    add_body_lines(slide, body_lines, Inches(0.75), Inches(1.25), Inches(7.65), Inches(4.6), 22)

    # Right-side story/visual box if there is supporting guidance.
    if visual_plan or discussion:
        panel = slide.shapes.add_shape(1, Inches(8.8), Inches(1.2), Inches(3.8), Inches(4.95))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(242, 246, 250)
        panel.line.color.rgb = RGBColor(185, 205, 225)
        if visual_plan:
            add_textbox(slide, "Visual / evidence plan", Inches(9.05), Inches(1.45), Inches(3.35), Inches(0.28), 13, True, (31, 78, 121))
            add_textbox(slide, visual_plan, Inches(9.05), Inches(1.78), Inches(3.35), Inches(1.9), 13, False, (40, 40, 40))
        if discussion:
            add_textbox(slide, "Audience prompt", Inches(9.05), Inches(3.85), Inches(3.35), Inches(0.28), 13, True, (31, 78, 121))
            add_textbox(slide, discussion, Inches(9.05), Inches(4.18), Inches(3.35), Inches(1.55), 13, False, (40, 40, 40))


def render_title_slide(prs: Presentation, deck: Dict[str, Any], slide_data: Dict[str, Any]) -> None:
    meta = deck["metadata"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    title = meta.get("presentation_title") or slide_data.get("title") or "Untitled Presentation"
    presenter = meta.get("presenter", "")
    date = meta.get("session_date", "")
    audience = meta.get("audience", "")
    talk_type = meta.get("presentation_type", "")

    add_textbox(slide, title, Inches(0.75), Inches(1.25), Inches(11.8), Inches(1.2), 34, True, (31, 78, 121))
    info_lines = [line for line in [presenter, date, audience, talk_type] if line]
    add_textbox(slide, "\n".join(info_lines), Inches(0.8), Inches(2.85), Inches(8.8), Inches(1.1), 18, False, (65, 65, 65))
    core_question = meta.get("core_question", "")
    if core_question:
        panel = slide.shapes.add_shape(1, Inches(0.8), Inches(4.5), Inches(11.6), Inches(1.05))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(234, 242, 250)
        panel.line.color.rgb = RGBColor(185, 205, 225)
        add_textbox(slide, "Core question", Inches(1.05), Inches(4.65), Inches(2.2), Inches(0.28), 13, True, (31, 78, 121))
        add_textbox(slide, core_question, Inches(1.05), Inches(4.95), Inches(10.9), Inches(0.38), 17, False, (45, 45, 45))


def render_objectives_slide(prs: Presentation, deck: Dict[str, Any], slide_data: Dict[str, Any], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, slide_output_title(deck, slide_data, index))
    objectives = split_nonempty_lines(slide_data.get("body", ""))
    if not objectives:
        objectives = OBJECTIVE_EXAMPLES[:3]
    add_body_lines(slide, objectives, Inches(0.85), Inches(1.35), Inches(11.5), Inches(4.8), 24)


def render_mentor_review_slide(prs: Presentation, deck: Dict[str, Any]) -> None:
    mr = deck["mentor_review"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Mentor review")
    lines = [
        f"Mentor: {mr.get('mentor_name') or 'Not listed'}",
        f"Review status: {mr.get('review_status') or 'Not sent'}",
    ]
    if mr.get("review_completed_date"):
        lines.append(f"Completed: {mr['review_completed_date']}")
    if mr.get("mentor_approval_statement"):
        lines.append(f"Approval statement: {mr['mentor_approval_statement']}")
    feedback = split_nonempty_lines(mr.get("mentor_feedback", ""))
    if feedback:
        lines.append("")
        lines.append("Mentor feedback summary:")
        lines.extend(feedback[:5])
    add_body_lines(slide, lines, Inches(0.85), Inches(1.3), Inches(11.7), Inches(5.0), 20)


def build_pptx(deck: Dict[str, Any]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    speaker_notes: List[str] = []
    output_index = 0
    for slide_data in deck["slides"]:
        if not slide_data.get("include", True):
            continue
        output_index += 1
        role = slide_data.get("role")
        if role == "Title":
            render_title_slide(prs, deck, slide_data)
        elif role == "Objectives":
            render_objectives_slide(prs, deck, slide_data, output_index)
        else:
            render_standard_slide(prs, deck, slide_data, output_index)
        speaker_notes.append(slide_data.get("speaker_notes", ""))

    if deck["mentor_review"].get("include_mentor_review_slide"):
        render_mentor_review_slide(prs, deck)
        speaker_notes.append("Internal mentor-review slide. Remove before final presentation if not intended for learners.")

    buffer = io.BytesIO()
    prs.save(buffer)
    pptx_bytes = buffer.getvalue()
    return add_speaker_notes_to_pptx(pptx_bytes, speaker_notes)


# -----------------------------------------------------------------------------
# Speaker notes injection for python-pptx output
# -----------------------------------------------------------------------------


def _next_rid(root: ET.Element) -> str:
    nums: List[int] = []
    for rel in root.findall(f"{{{REL_NS}}}Relationship"):
        rid = rel.get("Id", "")
        match = re.match(r"rId(\d+)$", rid)
        if match:
            nums.append(int(match.group(1)))
    return f"rId{max(nums or [0]) + 1}"


def _add_content_type_override(content_types_xml: bytes, part_name: str, content_type: str) -> bytes:
    ET.register_namespace("", CT_NS)
    root = ET.fromstring(content_types_xml)
    for override in root.findall(f"{{{CT_NS}}}Override"):
        if override.get("PartName") == part_name:
            return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    ET.SubElement(root, f"{{{CT_NS}}}Override", PartName=part_name, ContentType=content_type)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _notes_master_xml() -> bytes:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notesMaster xmlns:a="{A_NS}" xmlns:r="{R_NS}" xmlns:p="{P_NS}">
  <p:cSld>
    <p:bg><p:bgRef idx="1001"><a:schemeClr val="bg1"/></p:bgRef></p:bg>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg" idx="2"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="914400"/><a:ext cx="5486400" cy="3086100"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln w="12700"><a:solidFill><a:prstClr val="black"/></a:solidFill></a:ln></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
      </p:sp>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="4400550"/><a:ext cx="5486400" cy="3600450"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="en-US"/><a:t>Speaker notes</a:t></a:r></a:p></p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
</p:notesMaster>'''.encode("utf-8")


def _notes_slide_xml(notes: str, slide_number: int) -> bytes:
    lines = notes.splitlines() or [""]
    paragraphs = []
    for line in lines:
        escaped = html.escape(line, quote=False)
        paragraphs.append(
            f'<a:p><a:r><a:rPr lang="en-US" dirty="0"/><a:t>{escaped}</a:t></a:r><a:endParaRPr lang="en-US" dirty="0"/></a:p>'
        )
    paragraphs_xml = "".join(paragraphs)
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="{A_NS}" xmlns:r="{R_NS}" xmlns:p="{P_NS}">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr><p:spPr/></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/>{paragraphs_xml}</p:txBody></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="4" name="Slide Number Placeholder 3"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldNum" sz="quarter" idx="10"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:fld id="{{F7021451-1387-4CA6-816F-3879F97B5CBC}}" type="slidenum"><a:rPr lang="en-US"/><a:t>{slide_number}</a:t></a:fld><a:endParaRPr lang="en-US"/></a:p></p:txBody></p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>'''.encode("utf-8")


def _notes_slide_rels_xml(slide_number: int) -> bytes:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL_NS}">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{slide_number}.xml"/>
</Relationships>'''.encode("utf-8")


def add_speaker_notes_to_pptx(pptx_bytes: bytes, notes_by_slide: List[str]) -> bytes:
    """Add real PowerPoint speaker notes to a python-pptx-generated deck.

    python-pptx does not expose speaker-note authoring. This function post-processes
    the pptx package and adds notesSlide parts linked to each slide.
    """
    if not any((note or "").strip() for note in notes_by_slide):
        return pptx_bytes

    with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    # Content type for notes master.
    files["[Content_Types].xml"] = _add_content_type_override(
        files["[Content_Types].xml"],
        "/ppt/notesMasters/notesMaster1.xml",
        "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml",
    )

    # Add notes master relationship to presentation.
    presentation_rels_path = "ppt/_rels/presentation.xml.rels"
    pres_rels = ET.fromstring(files[presentation_rels_path])
    notes_master_rid = None
    for rel in pres_rels.findall(f"{{{REL_NS}}}Relationship"):
        if rel.get("Type", "").endswith("/notesMaster"):
            notes_master_rid = rel.get("Id")
            break
    if not notes_master_rid:
        notes_master_rid = _next_rid(pres_rels)
        ET.SubElement(
            pres_rels,
            f"{{{REL_NS}}}Relationship",
            Id=notes_master_rid,
            Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster",
            Target="notesMasters/notesMaster1.xml",
        )
        files[presentation_rels_path] = ET.tostring(pres_rels, encoding="utf-8", xml_declaration=True)

    # Add notes master ID list to presentation.xml if needed.
    presentation_path = "ppt/presentation.xml"
    pres = ET.fromstring(files[presentation_path])
    ns = {"p": P_NS, "r": R_NS}
    if pres.find("p:notesMasterIdLst", ns) is None:
        notes_master_id_list = ET.Element(f"{{{P_NS}}}notesMasterIdLst")
        notes_master_id = ET.SubElement(notes_master_id_list, f"{{{P_NS}}}notesMasterId")
        notes_master_id.set(f"{{{R_NS}}}id", notes_master_rid)
        insert_at = len(list(pres))
        for i, child in enumerate(list(pres)):
            if child.tag == f"{{{P_NS}}}sldSz":
                insert_at = i
                break
        pres.insert(insert_at, notes_master_id_list)
        files[presentation_path] = ET.tostring(pres, encoding="utf-8", xml_declaration=True)

    files["ppt/notesMasters/notesMaster1.xml"] = _notes_master_xml()
    files["ppt/notesMasters/_rels/notesMaster1.xml.rels"] = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL_NS}"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/></Relationships>'''.encode("utf-8")

    for slide_number, note in enumerate(notes_by_slide, start=1):
        if not (note or "").strip():
            continue
        slide_rels_path = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
        if slide_rels_path not in files:
            continue
        slide_rels = ET.fromstring(files[slide_rels_path])
        already_linked = any(rel.get("Type", "").endswith("/notesSlide") for rel in slide_rels.findall(f"{{{REL_NS}}}Relationship"))
        if not already_linked:
            rid = _next_rid(slide_rels)
            ET.SubElement(
                slide_rels,
                f"{{{REL_NS}}}Relationship",
                Id=rid,
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
                Target=f"../notesSlides/notesSlide{slide_number}.xml",
            )
            files[slide_rels_path] = ET.tostring(slide_rels, encoding="utf-8", xml_declaration=True)

        files[f"ppt/notesSlides/notesSlide{slide_number}.xml"] = _notes_slide_xml(note, slide_number)
        files[f"ppt/notesSlides/_rels/notesSlide{slide_number}.xml.rels"] = _notes_slide_rels_xml(slide_number)
        files["[Content_Types].xml"] = _add_content_type_override(
            files["[Content_Types].xml"],
            f"/ppt/notesSlides/notesSlide{slide_number}.xml",
            "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
        )

    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)
    return output.getvalue()


# -----------------------------------------------------------------------------
# Streamlit UI
# -----------------------------------------------------------------------------


def inject_css() -> None:
    st.markdown(
        """
        <style>
        .identity-card {
            border: 1px solid #d9e6f2;
            background: #f7fbff;
            border-radius: 12px;
            padding: 0.75rem 1rem;
            margin-bottom: 1rem;
        }
        .small-muted { color: #5c6b78; font-size: 0.9rem; }
        .slide-box {
            border: 1px solid #e5e7eb;
            border-radius: 14px;
            padding: 1rem;
            background: #ffffff;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_identity_card(deck: Dict[str, Any]) -> None:
    meta = deck["metadata"]
    title = meta.get("presentation_title") or "Untitled presentation"
    presenter = meta.get("presenter") or "Presenter not entered"
    date = meta.get("session_date") or "Date not entered"
    audience = meta.get("audience") or "Audience not entered"
    st.markdown(
        f"""
        <div class="identity-card">
        <strong>{html.escape(title)}</strong><br>
        <span class="small-muted">{html.escape(presenter)} · {html.escape(date)} · {html.escape(audience)}</span>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_sidebar(deck: Dict[str, Any]) -> None:
    st.sidebar.title("Slides")
    for idx, slide in enumerate(deck["slides"], start=1):
        marker = "✓" if slide.get("include", True) else "–"
        label_title = slide.get("title") or slide.get("role") or f"Slide {idx}"
        label = f"{idx}. {marker} {short_label(label_title)}"
        if st.sidebar.button(label, key=f"nav_{slide['id']}", use_container_width=True):
            st.session_state.selected_slide_id = slide["id"]

    st.sidebar.divider()
    if st.sidebar.button("+ Add story slide", use_container_width=True):
        new = new_slide(
            "Extra story slide",
            "",
            "Use this when the story needs another step, even if you do not know the title yet.",
            "Story",
            False,
        )
        # Insert before take-home points.
        insert_at = max(0, len(deck["slides"]) - 1)
        deck["slides"].insert(insert_at, new)
        st.session_state.selected_slide_id = new["id"]
        st.rerun()

    if st.sidebar.button("Reset draft", use_container_width=True):
        st.session_state.deck = default_deck()
        st.session_state.selected_slide_id = st.session_state.deck["slides"][0]["id"]
        st.rerun()


def render_metadata(deck: Dict[str, Any]) -> None:
    meta = deck["metadata"]
    with st.expander("Presentation identity", expanded=True):
        col1, col2 = st.columns([2, 1])
        with col1:
            meta["presentation_title"] = st.text_input("Presentation title", meta.get("presentation_title", ""))
            meta["presenter"] = st.text_input("Presenter", meta.get("presenter", ""))
        with col2:
            meta["session_date"] = st.text_input("Session date", meta.get("session_date", dt.date.today().isoformat()))
            meta["audience"] = st.text_input("Audience", meta.get("audience", "Pediatric residents"))

        col3, col4 = st.columns([1, 2])
        with col3:
            meta["presentation_type"] = st.selectbox(
                "Presentation type",
                ["Educational Topic", "Study Review", "Case Conference", "Hybrid Case + Evidence", "Research Update"],
                index=["Educational Topic", "Study Review", "Case Conference", "Hybrid Case + Evidence", "Research Update"].index(meta.get("presentation_type", "Educational Topic"))
                if meta.get("presentation_type", "Educational Topic") in ["Educational Topic", "Study Review", "Case Conference", "Hybrid Case + Evidence", "Research Update"]
                else 0,
            )
        with col4:
            meta["core_question"] = st.text_input(
                "Core question / narrative tension",
                meta.get("core_question", ""),
                help="Example: What should we do when the evidence and bedside workflow point in different directions?",
            )
        meta["story_arc"] = st.text_area(
            "Story arc",
            meta.get("story_arc", ""),
            height=90,
            help="One paragraph describing how the talk moves from problem → evidence/concepts → application → take-home points.",
        )


def render_bloom_helper() -> None:
    with st.expander("Bloom's taxonomy helper", expanded=True):
        st.write("Use verbs that describe what the learner will be able to do by the end of the session.")
        cols = st.columns(3)
        for idx, (level, verbs) in enumerate(BLOOM_HELPER.items()):
            with cols[idx % 3]:
                st.markdown(f"**{level}**")
                st.caption(verbs)
        st.markdown("**Examples**")
        for example in OBJECTIVE_EXAMPLES:
            st.caption(f"• {example}")


def render_slide_editor(deck: Dict[str, Any]) -> None:
    slide = get_selected_slide(deck)
    idx = deck["slides"].index(slide) + 1

    st.markdown(f"### Slide {idx}: {slide.get('role', 'Slide')}")
    st.caption(slide.get("prompt", ""))

    with st.container(border=True):
        col1, col2 = st.columns([2, 1])
        with col1:
            if slide.get("role") == "Title":
                st.info("The title slide pulls from Presentation identity. You can still add speaker notes here.")
            else:
                slide["title"] = st.text_input(
                    "Slide title",
                    slide.get("title", ""),
                    help="Leave blank if you do not know it yet. The exported deck will use a generic story-slide title.",
                )
                slide["subtitle"] = st.text_input("Optional subtitle / story signal", slide.get("subtitle", ""))
        with col2:
            slide["include"] = st.checkbox("Include in exported PPTX", bool(slide.get("include", True)), disabled=bool(slide.get("required", False)))
            if not slide.get("required", False):
                if st.button("Delete this slide", use_container_width=True):
                    deck["slides"].remove(slide)
                    st.session_state.selected_slide_id = deck["slides"][0]["id"]
                    st.rerun()

        if slide.get("role") == "Objectives":
            render_bloom_helper()
            slide["body"] = st.text_area(
                "Objectives — one per line",
                slide.get("body", ""),
                height=150,
                placeholder="Describe…\nDifferentiate…\nApply…",
            )
        elif slide.get("role") == "Disclosures":
            no_disclosures = st.checkbox("No relevant financial or non-financial disclosures", value="no relevant" in slide.get("body", "").lower())
            if no_disclosures and not slide.get("body", "").strip():
                slide["body"] = "No relevant financial or non-financial disclosures."
            slide["body"] = st.text_area("Disclosure text", slide.get("body", ""), height=110)
        elif slide.get("role") == "Title":
            st.write("Use the notes field below to script your opening hook.")
        else:
            slide["body"] = st.text_area(
                "Slide body — keep to one main idea",
                slide.get("body", ""),
                height=180,
                placeholder="Use short lines. Each line becomes a separate paragraph in the slide body.",
            )
            col3, col4 = st.columns(2)
            with col3:
                slide["visual_plan"] = st.text_area(
                    "Visual / evidence plan",
                    slide.get("visual_plan", ""),
                    height=110,
                    help="What figure, table, case data, graph, algorithm, or image belongs here?",
                )
            with col4:
                slide["discussion_prompt"] = st.text_area(
                    "Audience prompt",
                    slide.get("discussion_prompt", ""),
                    height=110,
                    help="Optional: question to ask residents or the room.",
                )

        slide["speaker_notes"] = st.text_area(
            "Speaker notes — exported into real PowerPoint speaker notes",
            slide.get("speaker_notes", ""),
            height=180,
            help="These notes are embedded as speaker notes in the exported PPTX file.",
        )


def render_mentor_review(deck: Dict[str, Any]) -> None:
    mr = deck["mentor_review"]
    with st.expander("Mentor review", expanded=False):
        col1, col2, col3 = st.columns([1.2, 1.2, 1])
        with col1:
            mr["mentor_name"] = st.text_input("Mentor name", mr.get("mentor_name", ""))
            mr["mentor_email"] = st.text_input("Mentor email", mr.get("mentor_email", ""))
        with col2:
            status_options = ["Not sent", "Sent for review", "Revisions requested", "Approved", "Final"]
            current_status = mr.get("review_status", "Not sent")
            mr["review_status"] = st.selectbox(
                "Review status",
                status_options,
                index=status_options.index(current_status) if current_status in status_options else 0,
            )
            mr["requested_review_date"] = st.text_input("Requested review date", mr.get("requested_review_date", ""))
        with col3:
            mr["review_completed_date"] = st.text_input("Review completed date", mr.get("review_completed_date", ""))
            mr["include_mentor_review_slide"] = st.checkbox("Include mentor-review slide in PPTX", bool(mr.get("include_mentor_review_slide", False)))

        st.markdown("**Mentor checklist**")
        for key, label in CHECKLIST_LABELS.items():
            mr["checklist"][key] = st.checkbox(label, bool(mr.get("checklist", {}).get(key, False)), key=f"mentor_{key}")

        mr["mentor_feedback"] = st.text_area("Mentor feedback / revision notes", mr.get("mentor_feedback", ""), height=130)
        mr["mentor_approval_statement"] = st.text_area(
            "Mentor approval statement",
            mr.get("mentor_approval_statement", ""),
            height=80,
            placeholder="Example: Reviewed with mentor; revisions completed; approved for presentation.",
        )


def render_story_check(deck: Dict[str, Any]) -> None:
    included = [slide for slide in deck["slides"] if slide.get("include", True)]
    missing_titles = [i + 1 for i, slide in enumerate(included) if not (slide.get("title") or "").strip() and slide.get("role") != "Title"]
    missing_notes = [i + 1 for i, slide in enumerate(included) if not (slide.get("speaker_notes") or "").strip()]
    objectives = next((s for s in deck["slides"] if s.get("role") == "Objectives"), None)
    objective_lines = split_nonempty_lines(objectives.get("body", "") if objectives else "")

    with st.expander("Story and readiness check", expanded=False):
        st.write("This is a lightweight quality-control check before exporting.")
        st.write(f"Included slides: **{len(included)}**")
        st.write(f"Objectives entered: **{len(objective_lines)}**")
        if missing_titles:
            st.warning(f"Slides without a title: {', '.join(map(str, missing_titles))}. This is allowed, but the deck will use generic story-slide titles.")
        else:
            st.success("All included non-title slides have titles.")
        if missing_notes:
            st.warning(f"Slides missing speaker notes: {', '.join(map(str, missing_notes))}.")
        else:
            st.success("All included slides have speaker notes.")


def render_archive_controls(deck: Dict[str, Any]) -> None:
    st.markdown("### Export / archive")
    pptx_bytes = build_pptx(deck)
    json_bytes = to_json_bytes(deck)
    slug = make_archive_slug(deck)
    pptx_filename = f"{slug}.pptx"
    json_filename = f"{slug}.json"

    col1, col2 = st.columns(2)
    with col1:
        st.download_button(
            "Download PowerPoint",
            data=pptx_bytes,
            file_name=pptx_filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
    with col2:
        st.download_button(
            "Download editable JSON draft",
            data=json_bytes,
            file_name=json_filename,
            mime="application/json",
            use_container_width=True,
        )

    uploaded = st.file_uploader("Reload an editable JSON draft", type=["json"])
    if uploaded is not None:
        try:
            st.session_state.deck = load_deck_from_json(uploaded.getvalue())
            st.session_state.selected_slide_id = st.session_state.deck["slides"][0]["id"]
            st.success("Draft loaded.")
            st.rerun()
        except Exception as exc:
            st.error(f"Could not load JSON draft: {exc}")

    config = github_config_from_secrets()
    with st.expander("GitHub archive", expanded=False):
        if not config:
            st.info(
                "GitHub saving is not configured. Add GITHUB_TOKEN, GITHUB_REPO, GITHUB_BRANCH, and GITHUB_FOLDER to Streamlit secrets."
            )
            st.code(
                'GITHUB_TOKEN = "ghp_your_token"\nGITHUB_REPO = "your_username/your_repo"\nGITHUB_BRANCH = "main"\nGITHUB_FOLDER = "presentation_archive"',
                language="toml",
            )
            return

        st.caption(f"Configured repo: {config.repo} · branch: {config.branch} · folder: {config.folder}")
        if st.button("Save draft + PowerPoint to GitHub", use_container_width=True):
            try:
                folder = f"{config.folder}/{slug}"
                github_upsert_bytes(config, f"{folder}/{ARCHIVE_JSON_NAME}", json_bytes, f"Save presentation draft: {slug}")
                github_upsert_bytes(config, f"{folder}/{ARCHIVE_PPTX_NAME}", pptx_bytes, f"Save presentation PPTX: {slug}")
                st.success(f"Saved to GitHub: {folder}")
            except Exception as exc:
                st.error(f"GitHub save failed: {exc}")

        if st.button("Refresh GitHub draft list", use_container_width=True):
            try:
                st.session_state.github_drafts = github_list_json_drafts(config)
            except Exception as exc:
                st.error(f"Could not list GitHub drafts: {exc}")

        drafts = st.session_state.get("github_drafts", [])
        if drafts:
            draft_names = [draft["name"] for draft in drafts]
            selected_name = st.selectbox("Load draft from GitHub", draft_names)
            selected = next(d for d in drafts if d["name"] == selected_name)
            if st.button("Load selected GitHub draft", use_container_width=True):
                try:
                    st.session_state.deck = github_load_json(config, selected["path"])
                    st.session_state.selected_slide_id = st.session_state.deck["slides"][0]["id"]
                    st.success("GitHub draft loaded.")
                    st.rerun()
                except Exception as exc:
                    st.error(f"Could not load selected draft: {exc}")


def main() -> None:
    st.set_page_config(page_title=APP_TITLE, page_icon="🩺", layout="wide")
    initialize_state()
    inject_css()

    deck = st.session_state.deck
    render_sidebar(deck)

    st.title(APP_TITLE)
    st.caption("Build a standardized, story-driven presentation with mentor review, GitHub archiving, and real PowerPoint speaker notes.")
    render_identity_card(deck)

    editor_col, export_col = st.columns([2.2, 0.95], gap="large")
    with editor_col:
        render_metadata(deck)
        render_slide_editor(deck)
        render_mentor_review(deck)
        render_story_check(deck)

    with export_col:
        render_archive_controls(deck)


if __name__ == "__main__":
    main()
